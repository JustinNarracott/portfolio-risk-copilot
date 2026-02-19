"""Unit tests for PPTX generator (Issue #21)."""

from datetime import date
from pathlib import Path

import pytest
from pptx import Presentation

from src.ingestion.parser import parse_file
from src.risk_engine.engine import analyse_portfolio
from src.artefacts.docx_generator import BrandConfig
from src.artefacts.pptx_generator import generate_board_slides

SAMPLE_CSV = Path(__file__).parent.parent.parent / "sample-data" / "jira-export-sample.csv"
REF_DATE = date(2026, 2, 19)


@pytest.fixture()
def report():
    projects = parse_file(SAMPLE_CSV)
    return analyse_portfolio(projects, top_n=5, reference_date=REF_DATE)


class TestBoardSlides:

    def test_generates_valid_pptx(self, report, tmp_path):
        out = tmp_path / "board.pptx"
        result = generate_board_slides(report, output_path=out)
        assert result.exists()
        prs = Presentation(str(result))
        assert len(prs.slides) == 2

    def test_slide1_has_title(self, report, tmp_path):
        out = tmp_path / "board.pptx"
        generate_board_slides(report, output_path=out)
        prs = Presentation(str(out))
        slide = prs.slides[0]
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        full = " ".join(texts)
        assert "Board Briefing" in full

    def test_slide1_has_project_names(self, report, tmp_path):
        out = tmp_path / "board.pptx"
        generate_board_slides(report, output_path=out)
        prs = Presentation(str(out))
        slide = prs.slides[0]
        texts = " ".join(shape.text_frame.text for shape in slide.shapes if shape.has_text_frame)
        # Should have at least some project names
        found = sum(1 for s in report.project_summaries if s.project_name in texts)
        assert found >= 3

    def test_slide2_has_risks(self, report, tmp_path):
        out = tmp_path / "board.pptx"
        generate_board_slides(report, output_path=out)
        prs = Presentation(str(out))
        slide = prs.slides[1]
        texts = " ".join(shape.text_frame.text for shape in slide.shapes if shape.has_text_frame)
        assert "Risk" in texts or "Critical" in texts or "High" in texts

    def test_custom_brand(self, report, tmp_path):
        brand = BrandConfig(primary_colour="990000", heading_font="Arial")
        out = tmp_path / "custom.pptx"
        result = generate_board_slides(report, brand=brand, output_path=out)
        assert result.exists()

    def test_default_output_path(self, report, tmp_path, monkeypatch):
        monkeypatch.chdir(tmp_path)
        result = generate_board_slides(report)
        assert result.name == "board-briefing.pptx"
