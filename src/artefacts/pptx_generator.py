"""
PPTX slide generator — professional board briefing slides.

Generates polished PowerPoint with charts, visual RAG indicators,
project cards, and data-driven risk summaries.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

from src.artefacts.docx_generator import BrandConfig
from src.risk_engine.engine import PortfolioRiskReport, Risk, RiskSeverity, RiskCategory

RAG_PPTX = {"Red": RGBColor(0xC0, 0x39, 0x2B), "Amber": RGBColor(0xE6, 0x7E, 0x22), "Green": RGBColor(0x27, 0xAE, 0x60)}
SEV_PPTX = {
    RiskSeverity.CRITICAL: RGBColor(0xC0, 0x39, 0x2B),
    RiskSeverity.HIGH: RGBColor(0xE6, 0x7E, 0x22),
    RiskSeverity.MEDIUM: RGBColor(0xF3, 0x9C, 0x12),
    RiskSeverity.LOW: RGBColor(0x27, 0xAE, 0x60),
}


def generate_board_slides(
    report: PortfolioRiskReport, brand: BrandConfig | None = None, output_path: str | Path | None = None,
) -> Path:
    """Generate board briefing PPTX with 2 slides."""
    brand = brand or BrandConfig()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    primary = RGBColor.from_string(brand.primary_colour)
    accent = RGBColor.from_string(brand.accent_colour)

    # ── Slide 1: Portfolio Dashboard ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Title
    _text(slide1, "Portfolio Health — Board Briefing",
          0.7, 0.3, 12, 0.7, size=Pt(28), bold=True,
          colour=primary, font=brand.heading_font)

    # Subtitle
    _text(slide1, f"{len(report.project_summaries)} projects  •  "
          f"{report.projects_at_risk} at risk  •  {report.total_risks} risks",
          0.7, 1.05, 8, 0.4, size=Pt(14),
          colour=RGBColor(0x70, 0x70, 0x70), font=brand.body_font)

    # Big RAG badge
    rag = report.portfolio_rag
    rag_col = RAG_PPTX.get(rag, RGBColor(0x80, 0x80, 0x80))
    badge = slide1.shapes.add_shape(1, Inches(10.5), Inches(0.4), Inches(2.2), Inches(1.0))
    badge.fill.solid()
    badge.fill.fore_color.rgb = rag_col
    badge.line.fill.background()
    _text(slide1, rag.upper(), 10.5, 0.45, 2.2, 0.55,
          size=Pt(24), bold=True, colour=RGBColor(0xFF, 0xFF, 0xFF),
          font=brand.heading_font, align=PP_ALIGN.CENTER)
    _text(slide1, "PORTFOLIO", 10.5, 0.95, 2.2, 0.3,
          size=Pt(9), bold=True, colour=RGBColor(0xFF, 0xFF, 0xFF),
          font=brand.body_font, align=PP_ALIGN.CENTER)

    # Project cards row
    card_y = 1.8
    card_w = 1.9
    card_h = 1.6
    gap = 0.15
    start_x = 0.7
    max_cards = min(len(report.project_summaries), 6)

    for i in range(max_cards):
        s = report.project_summaries[i]
        x = start_x + i * (card_w + gap)
        rc = RAG_PPTX.get(s.rag_status, RGBColor(0x80, 0x80, 0x80))

        # Card
        card = slide1.shapes.add_shape(1, Inches(x), Inches(card_y), Inches(card_w), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        card.line.fill.background()

        # Left accent bar
        accent_bar = slide1.shapes.add_shape(1, Inches(x), Inches(card_y), Inches(0.06), Inches(card_h))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = rc
        accent_bar.line.fill.background()

        _text(slide1, s.project_name, x + 0.15, card_y + 0.1, card_w - 0.25, 0.35,
              size=Pt(11), bold=True, colour=primary, font=brand.heading_font)
        _text(slide1, s.rag_status, x + 0.15, card_y + 0.5, card_w - 0.25, 0.35,
              size=Pt(18), bold=True, colour=rc, font=brand.heading_font)
        _text(slide1, f"{s.risk_count} risks  •  {s.project_status}",
              x + 0.15, card_y + 0.95, card_w - 0.25, 0.3,
              size=Pt(8), colour=RGBColor(0x70, 0x70, 0x70), font=brand.body_font)
        if s.risks:
            _text(slide1, s.risks[0].title[:50] + ("..." if len(s.risks[0].title) > 50 else ""),
                  x + 0.15, card_y + 1.2, card_w - 0.25, 0.3,
                  size=Pt(7), colour=RGBColor(0x95, 0xA5, 0xA6), font=brand.body_font)

    # RAG distribution chart
    reds = sum(1 for s in report.project_summaries if s.rag_status == "Red")
    ambers = sum(1 for s in report.project_summaries if s.rag_status == "Amber")
    greens = sum(1 for s in report.project_summaries if s.rag_status == "Green")

    if reds + ambers + greens > 0:
        chart_data = CategoryChartData()
        chart_data.categories = ["Red", "Amber", "Green"]
        chart_data.add_series("Projects", (reds, ambers, greens))
        chart_frame = slide1.shapes.add_chart(
            XL_CHART_TYPE.PIE, Inches(0.7), Inches(3.8), Inches(3.5), Inches(3.2), chart_data
        )
        chart = chart_frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.show_percentage = True
        plot.data_labels.show_value = False
        plot.data_labels.font.size = Pt(10)
        plot.data_labels.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # Colour the segments
        series = chart.series[0]
        for idx, colour in enumerate(["C0392B", "E67E22", "27AE60"]):
            pt = series.points[idx]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = RGBColor.from_string(colour)

    # Risk category breakdown chart
    cats: dict[str, int] = {}
    for s in report.project_summaries:
        for r in s.risks:
            cats[r.category.value] = cats.get(r.category.value, 0) + 1
    if cats:
        cat_data = CategoryChartData()
        sorted_cats = sorted(cats.items(), key=lambda x: -x[1])
        cat_data.categories = [c[0] for c in sorted_cats]
        cat_data.add_series("Count", [c[1] for c in sorted_cats])
        chart_frame2 = slide1.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, Inches(4.5), Inches(3.8), Inches(4.5), Inches(3.2), cat_data
        )
        chart2 = chart_frame2.chart
        chart2.has_legend = False
        plot2 = chart2.plots[0]
        plot2.has_data_labels = True
        plot2.data_labels.show_value = True
        plot2.data_labels.font.size = Pt(10)
        series2 = chart2.series[0]
        series2.format.fill.solid()
        series2.format.fill.fore_color.rgb = accent

    # Key decisions preview (right side)
    decisions = _get_decisions_text(report)
    _text(slide1, "KEY DECISIONS", 9.5, 3.8, 3.5, 0.35,
          size=Pt(11), bold=True, colour=primary,
          font=brand.heading_font)
    for di, dec in enumerate(decisions[:3]):
        short = dec[:90] + ("..." if len(dec) > 90 else "")
        _text(slide1, f"{di+1}. {short}", 9.5, 4.2 + di * 0.85, 3.5, 0.8,
              size=Pt(9), colour=RGBColor(0x50, 0x50, 0x50), font=brand.body_font)

    # ── Slide 2: Top Risks Detail ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)

    _text(slide2, "Top Portfolio Risks", 0.7, 0.3, 12, 0.6,
          size=Pt(24), bold=True, colour=primary, font=brand.heading_font)

    # Get top risks
    all_risks = []
    for s in report.project_summaries:
        all_risks.extend(s.risks)
    sev_order = {RiskSeverity.CRITICAL: 0, RiskSeverity.HIGH: 1, RiskSeverity.MEDIUM: 2, RiskSeverity.LOW: 3}
    all_risks.sort(key=lambda r: sev_order.get(r.severity, 99))
    top = all_risks[:4]

    # Risk cards
    for ri, risk in enumerate(top):
        cy = 1.2 + ri * 1.5
        sc = SEV_PPTX.get(risk.severity, RGBColor(0x80, 0x80, 0x80))

        # Card background
        card = slide2.shapes.add_shape(1, Inches(0.7), Inches(cy), Inches(11.9), Inches(1.3))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        card.line.fill.background()

        # Severity accent bar
        bar = slide2.shapes.add_shape(1, Inches(0.7), Inches(cy), Inches(0.06), Inches(1.3))
        bar.fill.solid()
        bar.fill.fore_color.rgb = sc
        bar.line.fill.background()

        # Severity badge
        _text(slide2, risk.severity.value.upper(), 0.95, cy + 0.1, 1.2, 0.3,
              size=Pt(10), bold=True, colour=sc, font=brand.heading_font)

        # Risk title
        _text(slide2, f"{risk.project_name}: {risk.title}",
              2.2, cy + 0.1, 10, 0.3,
              size=Pt(12), bold=True, colour=RGBColor(0x2C, 0x3E, 0x50), font=brand.heading_font)

        # Explanation (truncated)
        exp = risk.explanation[:220] + ("..." if len(risk.explanation) > 220 else "")
        _text(slide2, exp, 2.2, cy + 0.45, 10, 0.5,
              size=Pt(9), colour=RGBColor(0x60, 0x60, 0x60), font=brand.body_font)

        # Mitigation
        if risk.suggested_mitigation:
            mit = "→ " + risk.suggested_mitigation[:180] + ("..." if len(risk.suggested_mitigation) > 180 else "")
            _text(slide2, mit, 2.2, cy + 0.9, 10, 0.35,
                  size=Pt(8), colour=accent, font=brand.body_font, italic=True)

    # Save
    path = Path(output_path) if output_path else Path("board-briefing.pptx")
    prs.save(str(path))
    return path


# ──────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────

def _text(slide, text, left, top, width, height,
          size=Pt(12), bold=False, italic=False, colour=None,
          font="Calibri", align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if colour:
        run.font.color.rgb = colour


def _get_decisions_text(report) -> list[str]:
    """Quick decisions for slide preview."""
    decisions = []
    for s in report.project_summaries:
        for r in s.risks:
            if r.category == RiskCategory.BURN_RATE and r.severity == RiskSeverity.CRITICAL:
                decisions.append(f"{s.project_name}: Budget critical — approve top-up or cut scope")
                break
    blocked = set()
    for s in report.project_summaries:
        for r in s.risks:
            if r.category == RiskCategory.BLOCKED_WORK and r.severity in (RiskSeverity.CRITICAL, RiskSeverity.HIGH):
                blocked.add(s.project_name)
    if blocked:
        decisions.append(f"Unblock {', '.join(list(blocked)[:2])}: assign owners, 5-day deadline")
    reds = [s for s in report.project_summaries if s.rag_status == "Red"]
    if reds:
        decisions.append(f"Escalate {len(reds)} Red project{'s' if len(reds)>1 else ''} to executive review")
    if not decisions:
        decisions.append("Schedule portfolio risk review in 2 weeks")
    return decisions
